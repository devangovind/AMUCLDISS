import pptx
from pptx.util import Inches
import os
import re

class CreatePPT:
    def __init__(self):
        self.presentation = pptx.Presentation(os.path.join("fastapi-backend/pres", "templatelayout.pptx"))
        self.content_to_slide_dict = {}
        self.outputpath = os.path.join("pres", "output.pptx")
    
    def update_title(self, title):
        slide = self.presentation.slides[0]
        slide.shapes.title.text = title
        self.presentation.save(self.outputpath)

        

    def add_content(self, title, content):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        text_frame = slide.placeholders[1].text_frame
        text_frame.word_wrap = True
        self.presentation.save(self.outputpath)

    
    def list_text_boxes(self, slide_num):
        slide = self.presentation.slides[slide_num]
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                text_boxes.append(shape.text)
        return text_boxes
    def format_content(self, content):
        text = re.sub(r'### (.+)', r'\1', content)
        # Replace bold text denoted by '**'
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r"【.*?】", "", text)
        text = text.replace("\n", " ")
        return text


    def update_textbox(self, slide, to_replace, content):
        print("\n", content)
        try:
            slide = self.presentation.slides[self.content_to_slide_dict[slide]]
        except:
            return
        content = self.format_content(content)
        print("\n", content)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                if shape.text == to_replace:
                    text_frame = shape.text_frame
                    first_para = text_frame.paragraphs[0]
                    first_run = first_para.runs[0] if first_para.runs else first_para.add_run()
                    font = first_run.font
                    font_name  = font.name
                    font_size = font.size
                    # font_colour = font.color.rgb
                    text_frame.clear()
                    new_run = text_frame.paragraphs[0].add_run()
                    new_run.text = content
                    new_run.font.name = font_name
                    new_run.font.size = font_size 
                    # new_run.font.color.rgb = font_colour 
        # self.add_image_to_slide(slide, to_replace)
        self.presentation.save(self.outputpath)
    

    def add_image_to_slide(self,slide, metric):
        if metric == "kpis":
            slide = self.presentation.slides[5]
        else:
            try:
                slide = self.presentation.slides[self.content_to_slide_dict[slide]]
            except:
                return

        images = []
        if metric == "operatingincome": metric = "operating"
        if metric == "cashflow": metric = "cash"
        print("DIR", os.listdir("plots"))
        for file in os.listdir("plots"):
            print(file)
            if metric in file:
                images.append(os.path.join("plots", file))
            if metric == "kpis":
                if "operating" not in file and "cash" not in file and "revenue" not in file:
                    images.append(os.path.join("plots", file))

        print(images, metric)
        if metric != "kpis":
            x = 0
            for image in images:
                slide.shapes.add_picture(image, Inches(7.5), Inches(0.5+x), Inches(5))
                #                               left, top, width, height
                x+= 3.5
        else:
            x = 0
            row_counter = 0
            for image in images:
                slide.shapes.add_picture(image, Inches(0.5+x), Inches(1+((row_counter//2)*3.5)), Inches(5))
                #                               left, top, width, height
                row_counter += 1
                x+= 4

        self.presentation.save(self.outputpath)
        return 
if __name__ == "__main__":
    # print(os.path.join("pres", "input.pptx"))
    ppt = CreatePPT()
    # print(ppt.update_title("n"))
    ppt.add_content("test", "conent")