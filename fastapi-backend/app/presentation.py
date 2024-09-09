import pptx
from pptx.util import Inches, Pt
import os
import re

class PPT:
    def __init__(self):
        self.presentation = pptx.Presentation(os.path.join("pres", "templatelayout.pptx"))
        self.content_to_slide_dict = {}
        self.outputpath = os.path.join("pres", "output.pptx")
    
    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst  # pylint: disable=protected-access

    def move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    def update_title(self, title):
        slide = self.presentation.slides[0]
        slide.shapes.title.text = self.format_content(title)
        self.presentation.save(self.outputpath)

    def add_content(self, title, content):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = self.format_content(content)
        text_frame = slide.placeholders[1].text_frame
        text_frame.word_wrap = True

        self.presentation.save(self.outputpath)
    
    def add_images(self, metric):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = f"{metric} figures"
        images = []
        metric_key = metric.replace(" ", "").lower()
        for file in os.listdir("plots"):
            if metric_key in file:
                images.append(os.path.join("plots", file))
        column_counter = 0
        row_counter = 0
        for image in images:
            slide.shapes.add_picture(image, Inches(1+column_counter), Inches(1.5+((row_counter//2)*3)), Inches(5))
            #                               left,           top,                            width, 
            row_counter += 1
            column_counter+= 5.5
            if row_counter % 2  == 0:
                column_counter = 0
        new_index = self.find_slide(metric)
        self.move_slide(-1, new_index+1)
        self.presentation.save(self.outputpath)
        

    def find_slide(self, metric):
        i = 0
        for slide in self.presentation.slides:

            if slide.shapes.title.text == metric:
                return i
            i += 1
        return i

    def add_mda(self, scores):
        scores = scores.split(",")
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = "Sentiment Analysis Results"
        section = f"Overall: {scores[0]} \n Positive: {scores[1]} \n Neutral: {scores[2]} \n Negative: {scores[3]}"
        slide.placeholders[1].text = section
        self.presentation.save(self.outputpath)
    
    def format_content(self, content):
        text = re.sub(r'### (.+)', r'\1', content)
        # Replace bold text denoted by '**'
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r"【.*?】", "", text)
        text = text.replace("\n", " ")
        return text

